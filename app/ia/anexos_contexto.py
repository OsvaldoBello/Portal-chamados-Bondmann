"""Leitura de anexos (imagem/documento) para o contexto da triagem (F7).

Duas camadas: funções **puras** de processamento de bytes (testáveis sem
I/O — lição do incidente de 2026-07-29, quando um teste com conexão
totalmente mockada travou um valor que nunca passou pelo encoder real do
`asyncpg`) e um **orquestrador** assíncrono que baixa do Storage e nunca
lança (Regra de Ouro #5 — falha de anexo jamais derruba a triagem).

Escopo desta entrega: imagens (``image/jpeg``, ``image/png``) via visão e
documentos — PDF, Word (``.docx``), Excel (``.xlsx``) e PowerPoint
(``.pptx``) — via extração de texto. Vídeo e OCR de documento escaneado
ficam FORA de escopo — o dispatch por mime/extensão abaixo é a fronteira
estrutural, não só uma convenção de código.
"""

from __future__ import annotations

import asyncio
import base64
import io
import logging
from dataclasses import dataclass, field
from typing import Any

import openpyxl
from docx import Document
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from app.config import Settings
from app.storage import AnexosStorage

log = logging.getLogger("app.ia.anexos_contexto")

_MIME_IMAGEM = {"image/jpeg", "image/png"}
_MIME_PDF = {"application/pdf"}
_MIME_DOCX = {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
_MIME_XLSX = {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"}
_MIME_PPTX = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
_TIPOS_DOCUMENTO = {"pdf", "docx", "xlsx", "pptx"}


def _tipo_anexo(anexo: dict[str, Any]) -> str | None:
    """``"imagem"`` | ``"pdf"`` | ``"docx"`` | ``"xlsx"`` | ``"pptx"`` |
    ``None`` (fora do escopo desta leitura).

    docx/xlsx/pptx são contêineres ZIP (OOXML): libmagic antigo devolve
    ``application/zip`` em vez do mime específico — mesma ambiguidade já
    documentada em ``app/security/uploads.py``. Nesse caso, desempata pela
    extensão do ``path`` (atribuída pelo servidor no upload a partir da
    extensão validada — nunca o nome exibido, que é só metadado do
    cliente)."""
    mime = anexo.get("mime")
    if mime in _MIME_IMAGEM:
        return "imagem"
    if mime in _MIME_PDF:
        return "pdf"
    if mime in _MIME_DOCX:
        return "docx"
    if mime in _MIME_XLSX:
        return "xlsx"
    if mime in _MIME_PPTX:
        return "pptx"
    if mime == "application/zip":
        ext = (anexo.get("path") or "").rsplit(".", 1)[-1].lower()
        return ext if ext in {"docx", "xlsx", "pptx"} else None
    return None


def _truncar(texto: str, max_chars: int) -> str:
    texto = texto.strip()
    if not texto:
        return ""
    if len(texto) > max_chars:
        return texto[:max_chars].rstrip() + " […truncado]"
    return texto


def redimensionar_imagem(conteudo: bytes, *, max_dimensao: int, qualidade: int) -> str | None:
    """Decodifica, reduz ao maior lado ``max_dimensao`` (nunca amplia) e
    recomprime como JPEG — devolve uma *data URI* base64 pronta para o bloco
    ``image_url`` do modelo. ``None`` em bytes corrompidos/formato não
    suportado (nunca lança — quem chama decide a política de falha)."""
    try:
        imagem = Image.open(io.BytesIO(conteudo))
        imagem.load()  # força a decodificação agora (erros tardios do Pillow)
        imagem = imagem.convert("RGB")
        imagem.thumbnail((max_dimensao, max_dimensao), Image.LANCZOS)
        saida = io.BytesIO()
        imagem.save(saida, format="JPEG", quality=qualidade)
    except Exception:  # noqa: BLE001 — qualquer falha de decode/encode = ignora o anexo
        log.warning("Falha ao redimensionar imagem de anexo", exc_info=True)
        return None
    base64_str = base64.b64encode(saida.getvalue()).decode("ascii")
    return f"data:image/jpeg;base64,{base64_str}"


def extrair_texto_pdf(conteudo: bytes, *, max_paginas: int, max_chars: int) -> str:
    """Extrai texto das primeiras ``max_paginas`` páginas, truncado em
    ``max_chars``. ``""`` quando não há texto extraível (PDF escaneado/só
    imagem — sem OCR nesta entrega) ou o arquivo é inválido; nunca lança."""
    try:
        leitor = PdfReader(io.BytesIO(conteudo))
        partes = [
            (pagina.extract_text() or "").strip()
            for pagina in leitor.pages[:max_paginas]
        ]
    except Exception:  # noqa: BLE001 — PDF corrompido/criptografado = sem texto
        log.warning("Falha ao extrair texto de PDF anexado", exc_info=True)
        return ""
    return _truncar("\n\n".join(p for p in partes if p), max_chars)


def extrair_texto_docx(conteudo: bytes, *, max_chars: int) -> str:
    """Extrai texto de parágrafos e tabelas de um ``.docx``. ``""`` quando
    vazio/corrompido; nunca lança."""
    try:
        documento = Document(io.BytesIO(conteudo))
        partes = [p.text.strip() for p in documento.paragraphs if p.text.strip()]
        for tabela in documento.tables:
            for linha in tabela.rows:
                celulas = [c.text.strip() for c in linha.cells if c.text.strip()]
                if celulas:
                    partes.append(" | ".join(celulas))
    except Exception:  # noqa: BLE001 — docx corrompido = sem texto
        log.warning("Falha ao extrair texto de DOCX anexado", exc_info=True)
        return ""
    return _truncar("\n".join(partes), max_chars)


def extrair_texto_xlsx(conteudo: bytes, *, max_chars: int) -> str:
    """Extrai valores de célula de todas as planilhas em modo leitura
    otimizado (``read_only`` — não carrega o arquivo inteiro na memória) e
    para assim que atinge ``max_chars``, importante com o teto de 100MB do
    Químico. ``""`` quando vazio/corrompido; nunca lança."""
    try:
        pasta = openpyxl.load_workbook(io.BytesIO(conteudo), read_only=True, data_only=True)
    except Exception:  # noqa: BLE001 — xlsx corrompido = sem texto
        log.warning("Falha ao abrir XLSX anexado", exc_info=True)
        return ""
    partes: list[str] = []
    total = 0
    try:
        for planilha in pasta.worksheets:
            partes.append(f"## {planilha.title}")
            for linha in planilha.iter_rows(values_only=True):
                valores = [str(v).strip() for v in linha if v is not None and str(v).strip()]
                if not valores:
                    continue
                texto_linha = " | ".join(valores)
                partes.append(texto_linha)
                total += len(texto_linha)
                if total >= max_chars:
                    break
            if total >= max_chars:
                break
    except Exception:  # noqa: BLE001 — planilha malformada = sem texto
        log.warning("Falha ao ler células do XLSX anexado", exc_info=True)
        return ""
    finally:
        pasta.close()
    return _truncar("\n".join(partes), max_chars)


def extrair_texto_pptx(conteudo: bytes, *, max_chars: int) -> str:
    """Extrai texto das caixas de texto de cada slide de um ``.pptx``.
    ``""`` quando vazio/corrompido; nunca lança."""
    try:
        apresentacao = Presentation(io.BytesIO(conteudo))
        partes = []
        for indice, slide in enumerate(apresentacao.slides, start=1):
            textos = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if textos:
                partes.append(f"## Slide {indice}\n" + "\n".join(textos))
    except Exception:  # noqa: BLE001 — pptx corrompido = sem texto
        log.warning("Falha ao extrair texto de PPTX anexado", exc_info=True)
        return ""
    return _truncar("\n\n".join(partes), max_chars)


@dataclass(frozen=True)
class ResultadoAnexos:
    """Saída pronta para entrar em ambos os passes (Regra de Ouro #4: são
    arquivos que o próprio cliente subiu sobre o próprio problema — não são
    a base sigilosa da Bondmann — então reaproveitar entre Passe A e Passe B
    não vaza nada que o cliente já não tenha)."""

    blocos_imagem: list[dict[str, Any]] = field(default_factory=list)
    texto_documentos: str = ""
    processados: int = 0
    ignorados: int = 0


async def montar_contexto_anexos(
    storage: AnexosStorage,
    token: str,
    anexos: list[dict[str, Any]],
    settings: Settings,
    *,
    eh_quimico: bool,
) -> ResultadoAnexos:
    """Baixa e processa até os tetos configurados; cada item tem seu próprio
    isolamento de falha — download/decode ruim de um anexo nunca derruba os
    demais nem a triagem (Regra de Ouro #5)."""
    if not settings.ia_triagem_anexos_ativo or not anexos:
        return ResultadoAnexos()

    max_bytes = settings.ia_triagem_anexos_max_bytes_para(eh_quimico=eh_quimico)
    candidatos_imagem: list[tuple[dict[str, Any], str]] = []
    candidatos_documento: list[tuple[dict[str, Any], str]] = []
    ignorados = 0
    for anexo in anexos:
        tipo = _tipo_anexo(anexo)
        tamanho = anexo.get("tamanho") or 0
        if tipo == "imagem" and len(candidatos_imagem) < settings.ia_triagem_anexos_max_arquivos_imagem:
            alvo = candidatos_imagem
        elif tipo in _TIPOS_DOCUMENTO and len(candidatos_documento) < settings.ia_triagem_anexos_max_arquivos_documento:
            alvo = candidatos_documento
        else:
            ignorados += 1
            continue
        if tamanho > max_bytes:
            ignorados += 1
            continue
        alvo.append((anexo, tipo))

    candidatos = candidatos_imagem + candidatos_documento
    if not candidatos:
        return ResultadoAnexos(ignorados=ignorados)

    conteudos = await asyncio.gather(
        *(storage.download(token, anexo["path"]) for anexo, _tipo in candidatos),
        return_exceptions=True,
    )

    detail = settings.ia_triagem_anexos_detail_normalizado
    max_chars_doc = settings.ia_triagem_anexos_documento_max_chars
    blocos_imagem: list[dict[str, Any]] = []
    partes_documento: list[str] = []
    processados = 0
    for (anexo, tipo), conteudo in zip(candidatos, conteudos):
        if isinstance(conteudo, BaseException) or conteudo is None:
            if isinstance(conteudo, BaseException):
                log.warning("Falha ao baixar anexo %s: %s", anexo.get("path"), conteudo)
            ignorados += 1
            continue

        if tipo == "imagem":
            data_uri = redimensionar_imagem(
                conteudo,
                max_dimensao=settings.ia_triagem_anexos_max_dimensao_px,
                qualidade=settings.ia_triagem_anexos_qualidade_jpeg,
            )
            if data_uri is None:
                ignorados += 1
                continue
            blocos_imagem.append(
                {"type": "image_url", "image_url": {"url": data_uri, "detail": detail}}
            )
            processados += 1
            continue

        if tipo == "pdf":
            texto = extrair_texto_pdf(
                conteudo, max_paginas=settings.ia_triagem_anexos_pdf_max_paginas, max_chars=max_chars_doc
            )
        elif tipo == "docx":
            texto = extrair_texto_docx(conteudo, max_chars=max_chars_doc)
        elif tipo == "xlsx":
            texto = extrair_texto_xlsx(conteudo, max_chars=max_chars_doc)
        else:  # pptx
            texto = extrair_texto_pptx(conteudo, max_chars=max_chars_doc)

        if not texto:
            ignorados += 1
            continue
        nome = anexo.get("nome") or f"anexo.{tipo}"
        partes_documento.append(f"### {nome}\n{texto}")
        processados += 1

    return ResultadoAnexos(
        blocos_imagem=blocos_imagem,
        texto_documentos="\n\n".join(partes_documento),
        processados=processados,
        ignorados=ignorados,
    )
