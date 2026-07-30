"""Testes de `app/ia/anexos_contexto.py` (F7 — leitura de anexos pela IA).

Funções puras (`redimensionar_imagem`/`extrair_texto_pdf`) são testadas com
bytes reais via Pillow/pypdf — não só mocks. Essa escolha é deliberada: o
incidente de 2026-07-29 (`_MARGEM_ORFAO` como string em vez de `timedelta`)
só existiu porque o teste correspondente usava uma conexão 100% mockada e
nunca passou pelo encoder real do `asyncpg`. Aqui, o equivalente seria um
teste que nunca chama o Pillow/pypdf de verdade.
"""

from __future__ import annotations

import base64
import io

import openpyxl
from docx import Document
from PIL import Image
from pptx import Presentation

from app.config import Settings
from app.ia.anexos_contexto import (
    ResultadoAnexos,
    _tipo_anexo,
    extrair_texto_docx,
    extrair_texto_pdf,
    extrair_texto_pptx,
    extrair_texto_xlsx,
    montar_contexto_anexos,
    redimensionar_imagem,
)


def _imagem_png(largura: int, altura: int) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (largura, altura), color=(120, 40, 200)).save(buf, format="PNG")
    return buf.getvalue()


def _pdf_com_texto(texto: str) -> bytes:
    """PDF mínimo válido com uma página de texto real (sem dependências
    extras) — permite testar `extrair_texto_pdf` com bytes reais, iguais aos
    que o `pypdf` vê em produção."""
    conteudo = f"BT /F1 18 Tf 72 720 Td ({texto}) Tj ET".encode("latin-1")
    objetos = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> "
        b"/MediaBox [0 0 612 792] /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(conteudo)).encode() + b" >>\nstream\n" + conteudo + b"\nendstream",
    ]
    partes = [b"%PDF-1.4\n"]
    offsets = [0]
    for i, obj in enumerate(objetos, start=1):
        offsets.append(sum(len(p) for p in partes))
        partes.append(f"{i} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref_offset = sum(len(p) for p in partes)
    xref = [b"xref\n0 " + str(len(objetos) + 1).encode() + b"\n0000000000 65535 f \n"]
    for off in offsets[1:]:
        xref.append(f"{off:010} 00000 n \n".encode())
    trailer = (
        b"trailer\n<< /Size " + str(len(objetos) + 1).encode() + b" /Root 1 0 R >>\n"
        b"startxref\n" + str(xref_offset).encode() + b"\n%%EOF"
    )
    return b"".join(partes) + b"".join(xref) + trailer


def _pdf_sem_texto() -> bytes:
    """PDF válido, mas com página vazia (sem stream de conteúdo) — simula um
    PDF escaneado/só-imagem, que `extrair_texto_pdf` deve ignorar sem OCR."""
    objetos = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] >>",
    ]
    partes = [b"%PDF-1.4\n"]
    offsets = [0]
    for i, obj in enumerate(objetos, start=1):
        offsets.append(sum(len(p) for p in partes))
        partes.append(f"{i} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref_offset = sum(len(p) for p in partes)
    xref = [b"xref\n0 " + str(len(objetos) + 1).encode() + b"\n0000000000 65535 f \n"]
    for off in offsets[1:]:
        xref.append(f"{off:010} 00000 n \n".encode())
    trailer = (
        b"trailer\n<< /Size " + str(len(objetos) + 1).encode() + b" /Root 1 0 R >>\n"
        b"startxref\n" + str(xref_offset).encode() + b"\n%%EOF"
    )
    return b"".join(partes) + b"".join(xref) + trailer


def _docx_com_texto(paragrafo: str, *, tabela: list[list[str]] | None = None) -> bytes:
    documento = Document()
    documento.add_paragraph(paragrafo)
    if tabela:
        t = documento.add_table(rows=len(tabela), cols=len(tabela[0]))
        for i, linha in enumerate(tabela):
            for j, valor in enumerate(linha):
                t.cell(i, j).text = valor
    buf = io.BytesIO()
    documento.save(buf)
    return buf.getvalue()


def _xlsx_com_dados(linhas: list[list[str]], *, aba: str = "Planilha1") -> bytes:
    pasta = openpyxl.Workbook()
    ws = pasta.active
    ws.title = aba
    for linha in linhas:
        ws.append(linha)
    buf = io.BytesIO()
    pasta.save(buf)
    return buf.getvalue()


def _pptx_com_texto(*textos_por_slide: str) -> bytes:
    apresentacao = Presentation()
    layout = apresentacao.slide_layouts[5]  # layout em branco com título
    for texto in textos_por_slide:
        slide = apresentacao.slides.add_slide(layout)
        slide.shapes.title.text = texto
    buf = io.BytesIO()
    apresentacao.save(buf)
    return buf.getvalue()


def _settings(**overrides) -> Settings:
    base = dict(
        session_secret="segredo-real-de-teste-nao-default",
        csrf_secret="outro-segredo-real-de-teste-nao-default",
        ia_triagem_anexos_ativo=True,
    )
    base.update(overrides)
    return Settings(**base)


# --- redimensionar_imagem -----------------------------------------------


def test_redimensionar_imagem_produz_data_uri_jpeg_valida():
    uri = redimensionar_imagem(_imagem_png(2000, 1500), max_dimensao=1024, qualidade=78)
    assert uri.startswith("data:image/jpeg;base64,")
    decodificada = Image.open(io.BytesIO(base64.b64decode(uri.split(",", 1)[1])))
    assert decodificada.format == "JPEG"
    assert max(decodificada.size) <= 1024


def test_redimensionar_imagem_nao_amplia_imagem_pequena():
    uri = redimensionar_imagem(_imagem_png(200, 100), max_dimensao=1024, qualidade=78)
    decodificada = Image.open(io.BytesIO(base64.b64decode(uri.split(",", 1)[1])))
    assert decodificada.size == (200, 100)


def test_redimensionar_imagem_bytes_corrompidos_devolve_none():
    assert redimensionar_imagem(b"isto nao e uma imagem", max_dimensao=1024, qualidade=78) is None


# --- extrair_texto_pdf -----------------------------------------------


def test_extrair_texto_pdf_com_texto_real():
    texto = extrair_texto_pdf(
        _pdf_com_texto("BONDMANN TESTE EXTRACAO"), max_paginas=5, max_chars=6000
    )
    assert "BONDMANN TESTE EXTRACAO" in texto


def test_extrair_texto_pdf_trunca_no_max_chars():
    texto_longo = "A" * 100
    texto = extrair_texto_pdf(_pdf_com_texto(texto_longo), max_paginas=5, max_chars=20)
    assert len(texto) <= len("…truncado] ") + 20 + 5
    assert texto.endswith("[…truncado]")


def test_extrair_texto_pdf_sem_texto_extraivel_devolve_string_vazia():
    assert extrair_texto_pdf(_pdf_sem_texto(), max_paginas=5, max_chars=6000) == ""


def test_extrair_texto_pdf_bytes_corrompidos_devolve_string_vazia_sem_lancar():
    assert extrair_texto_pdf(b"isto nao e um pdf", max_paginas=5, max_chars=6000) == ""


# --- montar_contexto_anexos -----------------------------------------------


class _FakeStorage:
    def __init__(self, respostas: dict[str, bytes | None]):
        self._respostas = respostas
        self.chamadas: list[str] = []

    async def download(self, token, path):  # noqa: ARG002 — token não importa aqui
        self.chamadas.append(path)
        return self._respostas.get(path)


def _anexo(path: str, mime: str, tamanho: int = 1000, nome: str | None = None) -> dict:
    return {"path": path, "nome": nome or path, "mime": mime, "tamanho": tamanho}


async def test_kill_switch_off_nao_baixa_nada():
    storage = _FakeStorage({"a.jpg": _imagem_png(100, 100)})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.jpg", "image/jpeg")], _settings(ia_triagem_anexos_ativo=False),
        eh_quimico=False,
    )
    assert resultado == ResultadoAnexos()
    assert storage.chamadas == []


async def test_lista_de_anexos_vazia_nao_baixa_nada():
    storage = _FakeStorage({})
    resultado = await montar_contexto_anexos(storage, "token", [], _settings(), eh_quimico=False)
    assert resultado == ResultadoAnexos()
    assert storage.chamadas == []


async def test_mime_fora_do_allow_list_e_ignorado_sem_rede():
    storage = _FakeStorage({})
    resultado = await montar_contexto_anexos(
        storage, "token",
        [_anexo("v.mp4", "video/mp4"), _anexo("x.bin", "application/octet-stream")],
        _settings(), eh_quimico=False,
    )
    assert storage.chamadas == []
    assert resultado.processados == 0
    assert resultado.ignorados == 2


async def test_tamanho_acima_do_teto_e_ignorado_sem_baixar():
    settings = _settings(ia_triagem_anexos_max_bytes=1000)
    storage = _FakeStorage({"a.jpg": _imagem_png(100, 100)})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.jpg", "image/jpeg", tamanho=2000)], settings, eh_quimico=False,
    )
    assert storage.chamadas == []
    assert resultado.ignorados == 1


async def test_teto_maior_do_quimico_permite_arquivo_que_o_geral_rejeitaria():
    settings = _settings(ia_triagem_anexos_max_bytes=1000, ia_triagem_anexos_max_bytes_quimico=5000)
    conteudo = _imagem_png(100, 100)
    storage = _FakeStorage({"a.jpg": conteudo})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.jpg", "image/jpeg", tamanho=3000)], settings, eh_quimico=True,
    )
    assert storage.chamadas == ["a.jpg"]
    assert resultado.processados == 1


async def test_teto_por_tipo_limita_quantidade_de_downloads():
    settings = _settings(ia_triagem_anexos_max_arquivos_imagem=2)
    conteudo = _imagem_png(100, 100)
    anexos = [_anexo(f"img{i}.jpg", "image/jpeg") for i in range(5)]
    storage = _FakeStorage({a["path"]: conteudo for a in anexos})
    resultado = await montar_contexto_anexos(storage, "token", anexos, settings, eh_quimico=False)
    assert len(storage.chamadas) == 2
    assert resultado.processados == 2
    assert resultado.ignorados == 3


async def test_download_none_e_ignorado_sem_lancar():
    storage = _FakeStorage({"a.jpg": None})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.jpg", "image/jpeg")], _settings(), eh_quimico=False,
    )
    assert resultado.processados == 0
    assert resultado.ignorados == 1


async def test_imagem_corrompida_e_ignorada_sem_lancar():
    storage = _FakeStorage({"a.jpg": b"lixo binario"})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.jpg", "image/jpeg")], _settings(), eh_quimico=False,
    )
    assert resultado.processados == 0
    assert resultado.ignorados == 1


async def test_pdf_sem_texto_e_ignorado_sem_lancar():
    storage = _FakeStorage({"a.pdf": _pdf_sem_texto()})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.pdf", "application/pdf")], _settings(), eh_quimico=False,
    )
    assert resultado.processados == 0
    assert resultado.ignorados == 1


async def test_caso_feliz_imagem_e_pdf_juntos():
    storage = _FakeStorage(
        {"foto.jpg": _imagem_png(1200, 900), "laudo.pdf": _pdf_com_texto("RESULTADO DO LAUDO")}
    )
    resultado = await montar_contexto_anexos(
        storage, "token",
        [_anexo("foto.jpg", "image/jpeg", nome="foto.jpg"), _anexo("laudo.pdf", "application/pdf", nome="laudo.pdf")],
        _settings(), eh_quimico=False,
    )
    assert resultado.processados == 2
    assert resultado.ignorados == 0
    assert len(resultado.blocos_imagem) == 1
    assert resultado.blocos_imagem[0]["type"] == "image_url"
    assert resultado.blocos_imagem[0]["image_url"]["detail"] == "low"
    assert "RESULTADO DO LAUDO" in resultado.texto_documentos
    assert "laudo.pdf" in resultado.texto_documentos


async def test_detail_configuravel_e_valor_invalido_degrada_para_low():
    storage = _FakeStorage({"a.jpg": _imagem_png(100, 100)})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.jpg", "image/jpeg")],
        _settings(ia_triagem_anexos_detail="high"), eh_quimico=False,
    )
    assert resultado.blocos_imagem[0]["image_url"]["detail"] == "high"

    storage2 = _FakeStorage({"a.jpg": _imagem_png(100, 100)})
    resultado2 = await montar_contexto_anexos(
        storage2, "token", [_anexo("a.jpg", "image/jpeg")],
        _settings(ia_triagem_anexos_detail="algo-invalido"), eh_quimico=False,
    )
    assert resultado2.blocos_imagem[0]["image_url"]["detail"] == "low"


# --- extrair_texto_docx -----------------------------------------------

_MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def test_extrair_texto_docx_com_paragrafo_e_tabela():
    conteudo = _docx_com_texto("BONDMANN TESTE DOCX", tabela=[["Produto", "Lote"], ["X", "123"]])
    texto = extrair_texto_docx(conteudo, max_chars=6000)
    assert "BONDMANN TESTE DOCX" in texto
    assert "Produto | Lote" in texto
    assert "X | 123" in texto


def test_extrair_texto_docx_trunca_no_max_chars():
    texto = extrair_texto_docx(_docx_com_texto("A" * 100), max_chars=20)
    assert texto.endswith("[…truncado]")
    assert len(texto) <= 20 + len(" […truncado]")


def test_extrair_texto_docx_bytes_corrompidos_devolve_string_vazia_sem_lancar():
    assert extrair_texto_docx(b"isto nao e um docx", max_chars=6000) == ""


# --- extrair_texto_xlsx -----------------------------------------------


def test_extrair_texto_xlsx_com_dados_reais():
    conteudo = _xlsx_com_dados([["Produto", "Lote"], ["PRODTESTE X", "123"]], aba="Ocorrencias")
    texto = extrair_texto_xlsx(conteudo, max_chars=6000)
    assert "Ocorrencias" in texto
    assert "Produto | Lote" in texto
    assert "PRODTESTE X | 123" in texto


def test_extrair_texto_xlsx_para_de_ler_ao_bater_max_chars():
    # Sem o corte antecipado, isso produziria ~500*29 ≈ 14.500 caracteres.
    linhas = [[f"linha-{i}-" + "x" * 20] for i in range(500)]
    texto = extrair_texto_xlsx(_xlsx_com_dados(linhas), max_chars=100)
    assert texto.endswith("[…truncado]")
    assert len(texto) < 500  # parou bem antes de varrer as 500 linhas


def test_extrair_texto_xlsx_bytes_corrompidos_devolve_string_vazia_sem_lancar():
    assert extrair_texto_xlsx(b"isto nao e um xlsx", max_chars=6000) == ""


# --- extrair_texto_pptx -----------------------------------------------


def test_extrair_texto_pptx_com_texto_por_slide():
    conteudo = _pptx_com_texto("BONDMANN SLIDE UM", "BONDMANN SLIDE DOIS")
    texto = extrair_texto_pptx(conteudo, max_chars=6000)
    assert "Slide 1" in texto and "BONDMANN SLIDE UM" in texto
    assert "Slide 2" in texto and "BONDMANN SLIDE DOIS" in texto


def test_extrair_texto_pptx_trunca_no_max_chars():
    texto = extrair_texto_pptx(_pptx_com_texto("A" * 100), max_chars=20)
    assert texto.endswith("[…truncado]")


def test_extrair_texto_pptx_bytes_corrompidos_devolve_string_vazia_sem_lancar():
    assert extrair_texto_pptx(b"isto nao e um pptx", max_chars=6000) == ""


# --- _tipo_anexo (dispatch por mime + desempate por extensão) -------------


def test_tipo_anexo_reconhece_todos_os_mimes_especificos():
    assert _tipo_anexo(_anexo("a.jpg", "image/jpeg")) == "imagem"
    assert _tipo_anexo(_anexo("a.png", "image/png")) == "imagem"
    assert _tipo_anexo(_anexo("a.pdf", "application/pdf")) == "pdf"
    assert _tipo_anexo(_anexo("a.docx", _MIME_DOCX)) == "docx"
    assert _tipo_anexo(_anexo("a.xlsx", _MIME_XLSX)) == "xlsx"
    assert _tipo_anexo(_anexo("a.pptx", _MIME_PPTX)) == "pptx"


def test_tipo_anexo_desempata_application_zip_pela_extensao_do_path():
    """libmagic antigo devolve `application/zip` p/ OOXML (mesma ambiguidade
    documentada em app/security/uploads.py) — desempate pela extensão do
    `path` (atribuída pelo servidor, nunca o nome exibido)."""
    assert _tipo_anexo(_anexo("uuid123.docx", "application/zip")) == "docx"
    assert _tipo_anexo(_anexo("uuid123.xlsx", "application/zip")) == "xlsx"
    assert _tipo_anexo(_anexo("uuid123.pptx", "application/zip")) == "pptx"
    assert _tipo_anexo(_anexo("uuid123.zip", "application/zip")) is None  # zip de verdade


def test_tipo_anexo_mime_desconhecido_devolve_none():
    assert _tipo_anexo(_anexo("v.mp4", "video/mp4")) is None
    assert _tipo_anexo(_anexo("x.bin", "application/octet-stream")) is None


# --- montar_contexto_anexos: docx/xlsx/pptx via orquestrador ------------


async def test_orquestrador_processa_docx_xlsx_pptx_juntos():
    storage = _FakeStorage(
        {
            "a.docx": _docx_com_texto("TEXTO DO DOCX"),
            "a.xlsx": _xlsx_com_dados([["TEXTO DA XLSX"]]),
            "a.pptx": _pptx_com_texto("TEXTO DO PPTX"),
        }
    )
    resultado = await montar_contexto_anexos(
        storage, "token",
        [
            _anexo("a.docx", _MIME_DOCX, nome="a.docx"),
            _anexo("a.xlsx", _MIME_XLSX, nome="a.xlsx"),
            _anexo("a.pptx", _MIME_PPTX, nome="a.pptx"),
        ],
        _settings(ia_triagem_anexos_max_arquivos_documento=3), eh_quimico=False,
    )
    assert resultado.processados == 3
    assert resultado.ignorados == 0
    assert "TEXTO DO DOCX" in resultado.texto_documentos
    assert "TEXTO DA XLSX" in resultado.texto_documentos
    assert "TEXTO DO PPTX" in resultado.texto_documentos


async def test_orquestrador_reconhece_docx_com_mime_zip_ambiguo():
    storage = _FakeStorage({"a.docx": _docx_com_texto("TEXTO VIA ZIP AMBIGUO")})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.docx", "application/zip", nome="a.docx")],
        _settings(), eh_quimico=False,
    )
    assert resultado.processados == 1
    assert "TEXTO VIA ZIP AMBIGUO" in resultado.texto_documentos


async def test_teto_de_documento_e_compartilhado_entre_pdf_docx_xlsx_pptx():
    """`max_arquivos_documento` é um teto ÚNICO para os 4 tipos de documento
    (mesmo perfil de custo: texto extraído) — diferente do teto de imagem."""
    storage = _FakeStorage(
        {
            "a.pdf": _pdf_com_texto("PDF"),
            "b.docx": _docx_com_texto("DOCX"),
            "c.xlsx": _xlsx_com_dados([["XLSX"]]),
            "d.pptx": _pptx_com_texto("PPTX"),
        }
    )
    resultado = await montar_contexto_anexos(
        storage, "token",
        [
            _anexo("a.pdf", "application/pdf"),
            _anexo("b.docx", _MIME_DOCX),
            _anexo("c.xlsx", _MIME_XLSX),
            _anexo("d.pptx", _MIME_PPTX),
        ],
        _settings(ia_triagem_anexos_max_arquivos_documento=2), eh_quimico=False,
    )
    assert len(storage.chamadas) == 2
    assert resultado.processados == 2
    assert resultado.ignorados == 2


async def test_documento_sem_texto_e_ignorado_sem_lancar():
    storage = _FakeStorage({"a.docx": _docx_com_texto("")})
    resultado = await montar_contexto_anexos(
        storage, "token", [_anexo("a.docx", _MIME_DOCX)], _settings(), eh_quimico=False,
    )
    assert resultado.processados == 0
    assert resultado.ignorados == 1
